import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
import os

load_dotenv()

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

model = genai.GenerativeModel("gemini-1.5-flash")
response = model.generate_content("Explain how AI works")

print(response.text)


def create_presentation(slides_content, output_file="output_presentation.pptx"):

    # Create a presentation object
    prs = Presentation()

    for slide_data in slides_content:
        # Add a slide with title and content layout
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set slide title
        title.text = slide_data["title"]

        # Add bullet points
        for point in slide_data["content"]:
            p = content.text_frame.add_paragraph()
            p.text = point

    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as '{output_file}'")

# Example content for slides
slides_content = [
    {
        "title": "Introduction to Python",
        "content": [
            "Python is a high-level programming language.",
            "It is widely used for web development, data analysis, AI, and more.",
            "Python emphasizes readability and simplicity."
        ]
    },
    {
        "title": "Why Learn Python?",
        "content": [
            "Easy to learn and use.",
            "Large community and extensive libraries.",
            "Versatile across various fields like AI, web, and automation."
        ]
    },
    {
        "title": "Getting Started",
        "content": [
            "Install Python from the official website.",
            "Use an IDE like PyCharm or Visual Studio Code.",
            "Write your first Python script!"
        ]
    }
]

# Generate the presentation
create_presentation(slides_content)