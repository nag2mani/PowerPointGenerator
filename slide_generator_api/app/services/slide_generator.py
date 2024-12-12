from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import json
import os
import logging

def set_gradient_background(slide, color1, color2):
    """
    Apply a gradient background to a slide with two colors.
    """
    background = slide.background
    fill = background.fill
    fill.gradient()
    stop1 = fill.gradient_stops[0]
    stop2 = fill.gradient_stops[1]
    stop1.position, stop2.position = 0.0, 1.0
    stop1.color.rgb = RGBColor(*color1)
    stop2.color.rgb = RGBColor(*color2)

def generate_slides(content, presentation_id="presentation"):
    prs = Presentation()

    # Parse content if it is a JSON string
    if isinstance(content, str):
        try:
            content_json = json.loads(content)
        except json.JSONDecodeError as e:
            logging.error(f"Invalid JSON format: {e}")
            return None
    else:
        content_json = content

    for i, slide_content in enumerate(content_json):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using bullet layout
        title_text = slide_content.get("title", f"Slide {i+1}")
        content_text = slide_content.get("content", [])
        bg_color1 = slide_content.get("bg_color1", (255, 255, 255))  # Default white
        bg_color2 = slide_content.get("bg_color2", (0, 0, 0))        # Default black

        # Apply gradient background
        set_gradient_background(slide, bg_color1, bg_color2)

        # Title
        try:
            title = slide.shapes.title
            title.text = title_text
        except AttributeError:
            logging.warning(f"Slide {i+1} missing title placeholder.")

        # Content
        try:
            content = slide.placeholders[1]
            content.text_frame.clear()  # Clear default placeholder text
            if isinstance(content_text, str):
                content.text_frame.text = content_text
            elif isinstance(content_text, list):
                for point in content_text:
                    p = content.text_frame.add_paragraph()
                    p.text = point
                    p.font.size = Pt(18)  # Adjust font size if needed
            else:
                logging.warning(f"Unexpected content format on slide {i+1}.")
        except IndexError:
            logging.warning(f"Content placeholder missing on slide {i+1}.")

    presentations_dir = os.path.join(os.getcwd(), "presentations")
    os.makedirs(presentations_dir, exist_ok=True)
    output_path = os.path.join(presentations_dir, f"{presentation_id}.pptx")
    prs.save(output_path)
    logging.info(f"Presentation saved at '{output_path}'")
    return output_path