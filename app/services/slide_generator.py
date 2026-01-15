from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os
import logging
import re

# Beautiful color palettes for slides
COLOR_PALETTES = [
    # Blue to Purple gradient
    [(41, 128, 185), (142, 68, 173)],
    # Orange to Pink gradient
    [(255, 107, 107), (255, 159, 64)],
    # Teal to Blue gradient
    [(0, 206, 209), (30, 144, 255)],
    # Green to Teal gradient
    [(46, 213, 115), (0, 184, 148)],
    # Purple to Blue gradient
    [(108, 92, 231), (67, 97, 238)],
    # Red to Orange gradient
    [(255, 71, 87), (255, 165, 2)],
    # Indigo to Purple gradient
    [(99, 102, 241), (139, 92, 246)],
    # Cyan to Blue gradient
    [(6, 182, 212), (59, 130, 246)],
]

def parse_markdown_text(text):
    """
    Parse markdown-style formatting (**bold**, *italic*) and return segments with formatting info.
    Returns a list of tuples: (text, is_bold, is_italic)
    Handles cases like: "**Django:** A framework" or "text **bold** more text"
    """
    if not text:
        return [("", False, False)]
    
    segments = []
    i = 0
    
    while i < len(text):
        # Check for bold **text**
        if i < len(text) - 3 and text[i:i+2] == '**':
            # Find closing **
            end = text.find('**', i + 2)
            if end != -1:
                # Extract bold text
                bold_text = text[i+2:end]
                if bold_text:  # Only add if not empty
                    segments.append((bold_text, True, False))
                i = end + 2
                continue
        
        # Check for italic *text* (but not **)
        if i < len(text) - 1 and text[i] == '*' and (i == 0 or text[i-1] != '*') and (i == len(text) - 1 or text[i+1] != '*'):
            # Find closing *
            end = text.find('*', i + 1)
            if end != -1 and (end == len(text) - 1 or text[end+1] != '*'):  # Not part of **
                # Extract italic text
                italic_text = text[i+1:end]
                if italic_text:  # Only add if not empty
                    segments.append((italic_text, False, True))
                i = end + 1
                continue
        
        # Regular text - collect until we hit a formatting marker
        start = i
        while i < len(text):
            if text[i] == '*':
                # Check if it's start of ** or *
                if i < len(text) - 1 and text[i+1] == '*':
                    break  # Start of bold
                elif (i == 0 or text[i-1] != '*') and (i == len(text) - 1 or text[i+1] != '*'):
                    break  # Start of italic
            i += 1
        
        # Add regular text segment
        regular_text = text[start:i]
        if regular_text:
            segments.append((regular_text, False, False))
    
    return segments if segments else [(text, False, False)]

def add_formatted_text_to_paragraph(paragraph, text, font_size, font_color, font_name="Arial"):
    """
    Add text to a paragraph with markdown formatting support.
    Handles **bold** and *italic* formatting.
    """
    segments = parse_markdown_text(text)
    
    for idx, (segment_text, is_bold, is_italic) in enumerate(segments):
        if idx == 0:
            # Use the first run if paragraph already has text
            if paragraph.runs:
                run = paragraph.runs[0]
                run.text = segment_text
            else:
                run = paragraph.add_run()
                run.text = segment_text
        else:
            run = paragraph.add_run()
            run.text = segment_text
        
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.name = font_name
        run.font.bold = is_bold
        run.font.italic = is_italic

def set_gradient_background(slide, color1, color2):
    """
    Apply a beautiful gradient background to a slide with two colors.
    """
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0  # Diagonal gradient
    
    # Set gradient stops (python-pptx creates 2 stops by default)
    try:
        stop1 = fill.gradient_stops[0]
        stop1.position = 0.0
        stop1.color.rgb = RGBColor(*color1)
        
        # Ensure we have at least 2 stops
        if len(fill.gradient_stops) < 2:
            # Insert a second stop if it doesn't exist
            fill.gradient_stops.insert_gradient_stop(1.0, RGBColor(*color2))
        else:
            stop2 = fill.gradient_stops[1]
            stop2.position = 1.0
            stop2.color.rgb = RGBColor(*color2)
    except (IndexError, AttributeError):
        # Fallback: use solid color if gradient fails
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color1)

def generate_slides(content, presentation_id="presentation"):
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Parse content if it is a JSON string
    if isinstance(content, str):
        try:
            content_json = json.loads(content)
        except json.JSONDecodeError as e:
            logging.error(f"Invalid JSON format: {e}")
            return None
    else:
        content_json = content

    for i, slide_content in enumerate(content_json):
        # Use blank layout for more design control
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Select color palette (cycle through palettes)
        color_palette = COLOR_PALETTES[i % len(COLOR_PALETTES)]
        bg_color1, bg_color2 = color_palette
        
        # Apply beautiful gradient background
        set_gradient_background(slide, bg_color1, bg_color2)
        
        title_text = slide_content.get("title", f"Slide {i+1}")
        content_text = slide_content.get("content", [])
        
        # Define safe margins to keep everything inside
        MARGIN_LEFT = Inches(0.75)
        MARGIN_RIGHT = Inches(0.75)
        MARGIN_TOP = Inches(0.75)
        MARGIN_BOTTOM = Inches(0.75)
        
        # Calculate available width and height
        available_width = prs.slide_width - MARGIN_LEFT - MARGIN_RIGHT
        available_height = prs.slide_height - MARGIN_TOP - MARGIN_BOTTOM
        
        # Title section - fixed height
        TITLE_HEIGHT = Inches(1.2)
        TITLE_TOP = MARGIN_TOP + Inches(0.2)
        
        # Content section - remaining space
        CONTENT_TOP = TITLE_TOP + TITLE_HEIGHT + Inches(0.3)
        CONTENT_HEIGHT = available_height - TITLE_HEIGHT - Inches(0.5)
        
        try:
            # Add a semi-transparent rectangle for title background with proper margins
            title_bg_left = MARGIN_LEFT
            title_bg_top = TITLE_TOP
            title_bg_width = available_width
            title_bg_height = TITLE_HEIGHT
            
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                title_bg_left, title_bg_top, title_bg_width, title_bg_height
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
            title_bg.fill.transparency = 0.3  # Semi-transparent
            title_bg.line.fill.background()  # No border
            
            # Add title text box with proper padding inside the background
            title_padding = Inches(0.3)
            title_left = MARGIN_LEFT + title_padding
            title_top = TITLE_TOP + Inches(0.1)
            title_width = available_width - (title_padding * 2)
            title_height = TITLE_HEIGHT - Inches(0.2)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.word_wrap = True
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_frame.auto_size = None  # Disable auto-size to prevent overflow
            
            # Style the title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_run = title_paragraph.runs[0]
            title_run.font.size = Pt(38)  # Slightly smaller to fit better
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            title_run.font.name = "Arial"
        except Exception as e:
            logging.warning(f"Error adding title to slide {i+1}: {e}")

        # Add content area with proper boundaries
        try:
            content_padding = Inches(0.3)
            content_left = MARGIN_LEFT + content_padding
            content_top = CONTENT_TOP
            content_width = available_width - (content_padding * 2)
            content_height = CONTENT_HEIGHT - Inches(0.2)  # Leave some bottom margin
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.vertical_anchor = MSO_ANCHOR.TOP
            content_frame.auto_size = None  # Disable auto-size to prevent overflow
            content_frame.margin_left = Inches(0.2)
            content_frame.margin_right = Inches(0.2)
            content_frame.margin_top = Inches(0.1)
            content_frame.margin_bottom = Inches(0.1)
            
            if isinstance(content_text, str):
                # Single paragraph with markdown support
                p = content_frame.paragraphs[0]
                p.clear()  # Clear any default text
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)
                add_formatted_text_to_paragraph(
                    p, content_text, Pt(22), RGBColor(255, 255, 255)
                )
            elif isinstance(content_text, list):
                # Multiple bullet points with beautiful styling and markdown support
                for idx, point in enumerate(content_text):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                        p.text = ""  # Clear default text
                    else:
                        p = content_frame.add_paragraph()
                    
                    # Truncate very long points to prevent overflow
                    max_length = 150  # Maximum characters per bullet point
                    display_point = point[:max_length] + "..." if len(point) > max_length else point
                    
                    # Add bullet symbol (•) before text
                    bullet_run = p.add_run()
                    bullet_run.text = "•  "
                    bullet_run.font.size = Pt(24)
                    bullet_run.font.color.rgb = RGBColor(255, 255, 255)
                    bullet_run.font.name = "Arial"
                    bullet_run.font.bold = True
                    
                    # Parse and add formatted text with markdown support
                    segments = parse_markdown_text(display_point)
                    for seg_idx, (segment_text, is_bold, is_italic) in enumerate(segments):
                        text_run = p.add_run()
                        text_run.text = segment_text
                        text_run.font.size = Pt(20)
                        text_run.font.color.rgb = RGBColor(255, 255, 255)
                        text_run.font.name = "Arial"
                        text_run.font.bold = is_bold
                        text_run.font.italic = is_italic
                    
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(14)  # Reduced spacing
                    p.space_before = Pt(4)
                    p.level = 0
                    p.left_margin = Inches(0.1)
                    
            # Add a decorative accent shape within margins
            accent_left = MARGIN_LEFT + Inches(0.1)
            accent_top = CONTENT_TOP
            accent_width = Inches(0.12)
            accent_height = min(content_height, prs.slide_height - CONTENT_TOP - MARGIN_BOTTOM)
            
            accent = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, accent_left, accent_top, accent_width, accent_height
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(255, 255, 255)
            accent.fill.transparency = 0.2
            accent.line.fill.background()
            
            
        except Exception as e:
            logging.warning(f"Error adding content to slide {i+1}: {e}")

    presentations_dir = os.path.join(os.getcwd(), "presentations")
    os.makedirs(presentations_dir, exist_ok=True)
    output_path = os.path.join(presentations_dir, f"{presentation_id}.pptx")
    prs.save(output_path)
    logging.info(f"Presentation saved at '{output_path}'")
    return output_path